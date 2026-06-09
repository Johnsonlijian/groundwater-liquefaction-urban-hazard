"""Build editable Nature Water visual-upgrade figure pack.

Outputs:
- SVG/PDF/PNG for six main-figure upgrade candidates.
- Editable PPTX storyboard deck using PowerPoint shapes/text.

The generated figures are evidence-bounded scientific schematics. They do not
use AI-generated raster art as scientific evidence.
"""
from __future__ import annotations

import html
import json
import math
import textwrap
import argparse
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

try:
    import cairosvg  # type: ignore
except Exception:
    cairosvg = None


ROOT = Path(__file__).resolve().parents[1]
PROJECT = ROOT.parent
DATA = ROOT / "data"
OUT = ROOT / "export"
SVG_DIR = OUT / "svg"
PDF_DIR = OUT / "pdf"
PNG_DIR = OUT / "png_600dpi"
PPT_DIR = OUT / "pptx"
for d in [SVG_DIR, PDF_DIR, PNG_DIR, PPT_DIR]:
    d.mkdir(parents=True, exist_ok=True)

W, H = 2200, 1500
NAVY = "#081B4A"
BLUE = "#2467AE"
TEAL = "#2F9EB0"
SKY = "#79B7D8"
CORAL = "#C84C3A"
RED = "#B51F2B"
ORANGE = "#D86B20"
GREEN = "#2E7D3A"
LIGHT_GREEN = "#EAF5E8"
AMBER = "#F2B84B"
PALE = "#F7FAFE"
PANEL = "#FFFFFF"
BORDER = "#A9BCD6"
GREY = "#5E6A7D"
LIGHT_GREY = "#E8EDF4"
MID_GREY = "#B6C0CC"
BLACK = "#111827"


def esc(s: object) -> str:
    return html.escape(str(s), quote=True)


def wrap_lines(s: str, width: int) -> list[str]:
    return textwrap.wrap(str(s), width=width, break_long_words=False, replace_whitespace=False) or [""]


def svg_header(title: str, subtitle: str) -> list[str]:
    return [
        f'<svg xmlns="http://www.w3.org/2000/svg" width="{W}" height="{H}" viewBox="0 0 {W} {H}">',
        "<defs>",
        '<filter id="shadow" x="-20%" y="-20%" width="140%" height="140%"><feDropShadow dx="0" dy="8" stdDeviation="10" flood-color="#0b1b3a" flood-opacity="0.16"/></filter>',
        '<linearGradient id="water" x1="0" x2="1"><stop offset="0" stop-color="#8BD5F7" stop-opacity=".9"/><stop offset="1" stop-color="#1B77B9" stop-opacity=".9"/></linearGradient>',
        '<linearGradient id="soil" x1="0" x2="0" y1="0" y2="1"><stop offset="0" stop-color="#C28B59"/><stop offset=".55" stop-color="#B1784A"/><stop offset="1" stop-color="#7C4D2D"/></linearGradient>',
        '<marker id="arrow" viewBox="0 0 10 10" refX="9" refY="5" markerWidth="9" markerHeight="9" orient="auto-start-reverse"><path d="M0,0 L10,5 L0,10 z" fill="#081B4A"/></marker>',
        '<marker id="arrowBlue" viewBox="0 0 10 10" refX="9" refY="5" markerWidth="9" markerHeight="9" orient="auto-start-reverse"><path d="M0,0 L10,5 L0,10 z" fill="#2467AE"/></marker>',
        '<marker id="arrowRed" viewBox="0 0 10 10" refX="9" refY="5" markerWidth="9" markerHeight="9" orient="auto-start-reverse"><path d="M0,0 L10,5 L0,10 z" fill="#C84C3A"/></marker>',
        "</defs>",
        f'<rect x="0" y="0" width="{W}" height="{H}" rx="0" fill="#F5F8FC"/>',
        text(W / 2, 58, title.upper(), 46, NAVY, bold=True, anchor="middle"),
        text(W / 2, 102, subtitle, 25, NAVY, italic=True, anchor="middle"),
    ]


def svg_end() -> str:
    return "</svg>\n"


def text(x, y, s, size=24, color=BLACK, bold=False, italic=False, anchor="start", family="Arial", opacity=1.0):
    weight = "700" if bold else "400"
    style = "font-style:italic;" if italic else ""
    return (
        f'<text x="{x:.1f}" y="{y:.1f}" font-family="{family}" font-size="{size}" '
        f'font-weight="{weight}" text-anchor="{anchor}" fill="{color}" opacity="{opacity}" style="{style}">{esc(s)}</text>'
    )


def multiline(x, y, s, size=22, color=BLACK, width=42, line_h=1.18, bold=False, italic=False, anchor="start"):
    out = []
    for i, line in enumerate(wrap_lines(str(s), width)):
        out.append(text(x, y + i * size * line_h, line, size, color, bold, italic, anchor))
    return "\n".join(out)


def rect(x, y, w, h, fill=PANEL, stroke=BORDER, sw=2, rx=22, shadow=False, opacity=1):
    filt = ' filter="url(#shadow)"' if shadow else ""
    return f'<rect x="{x}" y="{y}" width="{w}" height="{h}" rx="{rx}" fill="{fill}" stroke="{stroke}" stroke-width="{sw}" opacity="{opacity}"{filt}/>'


def circle(x, y, r, fill=BLUE, stroke="white", sw=3):
    return f'<circle cx="{x}" cy="{y}" r="{r}" fill="{fill}" stroke="{stroke}" stroke-width="{sw}"/>'


def line(x1, y1, x2, y2, color=NAVY, sw=4, dash=None, marker=None, opacity=1):
    dash_attr = f' stroke-dasharray="{dash}"' if dash else ""
    marker_attr = f' marker-end="url(#{marker})"' if marker else ""
    return f'<line x1="{x1}" y1="{y1}" x2="{x2}" y2="{y2}" stroke="{color}" stroke-width="{sw}" stroke-linecap="round"{dash_attr}{marker_attr} opacity="{opacity}"/>'


def path(d, fill="none", stroke=NAVY, sw=3, marker=None, opacity=1):
    marker_attr = f' marker-end="url(#{marker})"' if marker else ""
    return f'<path d="{d}" fill="{fill}" stroke="{stroke}" stroke-width="{sw}" stroke-linecap="round" stroke-linejoin="round"{marker_attr} opacity="{opacity}"/>'


def polygon(points, fill, stroke=BORDER, sw=2, opacity=1):
    pts = " ".join(f"{x},{y}" for x, y in points)
    return f'<polygon points="{pts}" fill="{fill}" stroke="{stroke}" stroke-width="{sw}" opacity="{opacity}"/>'


def panel_label(x, y, letter, color=BLUE):
    return circle(x, y, 25, color) + text(x, y + 9, letter, 26, "white", True, anchor="middle")


def badge(x, y, label, fill, color="white", w=170):
    return rect(x, y, w, 40, fill, fill, 1.5, 20) + text(x + w / 2, y + 27, label, 17, color, True, anchor="middle")


def export_svg(name: str, parts: list[str]):
    svg = "\n".join(parts) + "\n" + svg_end()
    path_svg = SVG_DIR / f"{name}.svg"
    path_pdf = PDF_DIR / f"{name}.pdf"
    path_png = PNG_DIR / f"{name}.png"
    path_svg.write_text(svg, encoding="utf-8")
    if cairosvg is not None:
        try:
            cairosvg.svg2pdf(bytestring=svg.encode("utf-8"), write_to=str(path_pdf))
            cairosvg.svg2png(bytestring=svg.encode("utf-8"), write_to=str(path_png), output_width=W * 2, output_height=H * 2)
            return path_svg
        except Exception:
            pass

    # Fallback for Windows environments without the native Cairo DLL.
    # PyMuPDF opens SVG directly and keeps the SVG as the primary editable asset.
    try:
        import fitz

        doc = fitz.open(str(path_svg))
        path_pdf.write_bytes(doc.convert_to_pdf())
        pix = doc[0].get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
        pix.save(str(path_png))
        doc.close()
        return path_svg
    except Exception:
        pass

    from reportlab.graphics import renderPDF
    from svglib.svglib import svg2rlg
    import fitz

    drawing = svg2rlg(str(path_svg))
    renderPDF.drawToFile(drawing, str(path_pdf), showBoundary=0)
    doc = fitz.open(str(path_pdf))
    pix = doc[0].get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
    pix.save(str(path_png))
    doc.close()
    return path_svg


def load_summary():
    display = json.loads((DATA / "display_summary.json").read_text(encoding="utf-8"))
    readiness = pd.DataFrame()
    regional = pd.read_csv(DATA / "regional_groups.csv")
    ladder = pd.read_csv(DATA / "product_consensus.csv")
    cards = pd.read_csv(DATA / "evidence_cards.csv")
    phase = pd.read_csv(DATA / "sy_phase.csv")
    eng = pd.read_csv(DATA / "engineering_enrichment.csv")
    events = pd.read_csv(DATA / "event_boundary.csv")
    protocol = pd.read_csv(DATA / "protocol_steps.csv")
    cities = pd.read_csv(DATA / "dgls_city_outputs.csv")
    return display, readiness, regional, ladder, cards, phase, eng, events, protocol, cities


DISPLAY, READINESS, REGIONAL, LADDER, CARDS, PHASE, ENG, EVENTS, PROTOCOL, CITIES = load_summary()
COUNTS = DISPLAY["global_payload_counts"]


def fig1_mechanism():
    p = svg_header("Figure 1  |  Dynamic Groundwater Review Screen", "A 3D state-variable mechanism, not a site-scale hazard map")
    p += [rect(32, 128, 2164, 1270, "#FFFFFF", BORDER, 2.2, 28, True)]

    # 3D aquifer cutaway
    ox, oy = 190, 330
    top = [(ox + 120, oy), (ox + 720, oy + 70), (ox + 550, oy + 250), (ox - 60, oy + 170)]
    side = [(ox - 60, oy + 170), (ox + 550, oy + 250), (ox + 550, oy + 620), (ox - 60, oy + 540)]
    front = [(ox + 550, oy + 250), (ox + 720, oy + 70), (ox + 720, oy + 420), (ox + 550, oy + 620)]
    water = [(ox + 10, oy + 350), (ox + 580, oy + 410), (ox + 680, oy + 300), (ox + 90, oy + 240)]
    p += [
        panel_label(88, 185, "A", BLUE),
        text(126, 193, "3D aquifer-state mechanism", 26, NAVY, True),
        polygon(side, "url(#soil)", "#8A5A37", 3),
        polygon(front, "#9D6941", "#7A5032", 3),
        polygon(water, "url(#water)", "#1B77B9", 2.5, 0.78),
        polygon(top, "#7FA16C", "#5D764D", 3),
        path(f"M{ox+30},{oy+345} C{ox+210},{oy+305} {ox+365},{oy+405} {ox+590},{oy+395}", "none", "#DDF7FF", 8, opacity=.55),
        text(ox + 280, oy + 438, "regional water-table perturbation", 20, "#0E5D8F", True, anchor="middle"),
        text(ox + 282, oy + 464, "ΔWTD = storage trend / S_y", 19, "#0E5D8F", anchor="middle"),
    ]
    # City blocks
    for i, (bx, by, bh, col) in enumerate([(210, 350, 85, "#D9E5F3"), (290, 365, 130, "#C7D8EA"), (380, 350, 105, "#B7CDE2"), (510, 370, 150, "#D5DFEA")]):
        p += [rect(bx, by - bh, 55, bh, col, "#6A7A8C", 2, 4)]
        p += [polygon([(bx, by - bh), (bx + 28, by - bh - 22), (bx + 83, by - bh - 22), (bx + 55, by - bh)], "#EEF4FA", "#6A7A8C", 1.5)]
    # Soil grains / liquefaction bubbles
    for i in range(34):
        x = ox + 25 + (i * 61) % 520
        y = oy + 480 + ((i * 37) % 90)
        p.append(circle(x, y, 8 + (i % 3), "#D6B08A", "#7B5132", 1))
    for i in range(10):
        p.append(circle(ox + 430 + i * 18, oy + 430 - i * 9, 7, "#E9FCFF", "#65B8D9", 1.2))

    # Satellite and driver
    sx, sy = 1090, 275
    p += [
        panel_label(900, 185, "B", TEAL),
        text(940, 193, "GRACE regional storage", 25, NAVY, True),
        circle(sx, sy, 52, "#1F4D83", "#D6E7FF", 4),
        polygon([(sx - 145, sy - 20), (sx - 62, sy - 48), (sx - 40, sy - 18), (sx - 122, sy + 12)], "#D8EAFE", "#356CA8", 2),
        polygon([(sx + 62, sy - 48), (sx + 145, sy - 20), (sx + 122, sy + 12), (sx + 40, sy - 18)], "#D8EAFE", "#356CA8", 2),
        text(sx, sy + 10, "GRACE", 21, "white", True, anchor="middle"),
        path(f"M{sx-20},{sy+72} C{sx-190},{sy+170} {sx-360},{sy+200} {ox+650},{oy+150}", "none", TEAL, 5, "arrowBlue"),
        badge(935, 380, "300-km storage signal", TEAL, "white", 245),
        multiline(945, 455, "Regional driver; city is the exposure marker, not an independent groundwater well.", 19, GREY, 29),
    ]

    # Model term and screening dial
    mx, my = 1360, 565
    p += [
        panel_label(1255, 185, "C", CORAL),
        text(1295, 193, "Published water-table term", 25, NAVY, True),
        rect(1245, 245, 460, 510, "#FFF7F2", "#F0B09E", 2.2, 24, True),
        text(1475, 300, "Zhu et al. geospatial screen", 26, NAVY, True, anchor="middle"),
        text(1475, 350, "all non-groundwater predictors fixed", 19, GREY, anchor="middle"),
        line(1320, 500, 1620, 500, "#9AA8B8", 4),
        path("M1325,500 A150,150 0 0,1 1625,500", "none", "#9AA8B8", 5),
        path("M1475,500 L1578,415", "none", CORAL, 8),
        circle(1475, 500, 12, CORAL, "white", 3),
        text(1475, 610, "water-table coefficient only", 25, CORAL, True, anchor="middle"),
        text(1475, 648, "ΔP_liq screening increment", 23, BLACK, anchor="middle"),
    ]
    p += [
        path(f"M{ox+735},{oy+320} C{970},{my-30} {1110},{my-30} {1245},{my}", "none", NAVY, 5, "arrow"),
        text(970, 670, "specific-yield phase", 21, NAVY, True),
    ]

    # Review output protocol
    p += [
        panel_label(1748, 185, "D", GREEN),
        text(1788, 193, "Local review action", 25, NAVY, True),
        rect(1740, 245, 430, 510, "#F7FBF6", "#A8CFA3", 2.2, 24, True),
    ]
    actions = [
        ("A", "material follow-up", RED),
        ("B", "targeted local data", ORANGE),
        ("C", "detectable update", BLUE),
        ("D", "routine refresh", GREY),
    ]
    for i, (lab, desc, col) in enumerate(actions):
        y = 315 + i * 82
        p += [circle(1800, y, 26, col, "white", 3), text(1800, y + 8, lab, 23, "white", True, anchor="middle"), text(1842, y - 2, desc, 24, NAVY, True), multiline(1842, y + 28, "local wells + CPT/SPT + sediment review", 16, GREY, 38)]
    p.append(path("M1706,535 C1770,535 1775,535 1740,535", "none", GREEN, 5, "arrowBlue"))

    # Bottom takeaways
    p += [
        rect(70, 1130, 2050, 230, "#F6F9FD", "#B8C8DD", 2, 24),
        text(108, 1180, "KEY MECHANISM", 27, NAVY, True),
        badge(340, 1150, "recharge → shallower WTD → P_liq ↑", CORAL, "white", 430),
        badge(820, 1150, "depletion → deeper WTD → P_liq ↓", BLUE, "white", 430),
        badge(1300, 1150, "subsidence + water-security audit", AMBER, NAVY, 420),
        multiline(108, 1236, "The scientific contribution is a dynamic groundwater state-variable screen. It flags where a static water-table assumption becomes stale; it does not predict earthquakes, local damage, or engineering factor of safety.", 24, BLACK, 138),
    ]
    return export_svg("Fig1_3D_DGLS_mechanism_upgrade", p)


def fig2_global_payload():
    p = svg_header("Figure 2  |  Global Null, Regional Payload", "The effect is not diffuse global amplification; it is a regional review queue")
    p += [rect(32, 128, 2164, 1270, "#FFFFFF", BORDER, 2.2, 28, True)]
    p += [panel_label(85, 185, "A", BLUE), text(124, 193, "444 seismic city exposure markers", 29, NAVY, True)]
    # Map panel
    mx, my, mw, mh = 70, 240, 980, 620
    p += [rect(mx, my, mw, mh, "#F2F5F9", "#B9C7D9", 2, 24)]
    for lon in range(-120, 181, 60):
        x = mx + (lon + 180) / 360 * mw
        p.append(line(x, my + 20, x, my + mh - 20, "#D5DDE8", 1.4, "3 9"))
    for lat in range(-60, 61, 30):
        y = my + (90 - lat) / 180 * mh
        p.append(line(mx + 20, y, mx + mw - 20, y, "#D5DDE8", 1.4, "3 9"))
    # Simplified continents
    p += [
        polygon([(mx+120,my+170),(mx+300,my+90),(mx+390,my+175),(mx+320,my+285),(mx+185,my+270)], "#D7DEE8", "#C2CCD8", 1.3, .9),
        polygon([(mx+350,my+360),(mx+430,my+310),(mx+505,my+435),(mx+450,my+535),(mx+360,my+510)], "#D7DEE8", "#C2CCD8", 1.3, .9),
        polygon([(mx+505,my+165),(mx+650,my+95),(mx+800,my+135),(mx+885,my+250),(mx+740,my+330),(mx+590,my+275)], "#D7DEE8", "#C2CCD8", 1.3, .9),
        polygon([(mx+725,my+365),(mx+860,my+385),(mx+875,my+500),(mx+740,my+505)], "#D7DEE8", "#C2CCD8", 1.3, .9),
    ]
    # City dots actual lon/lat
    for _, r in CITIES.iterrows():
        x = mx + (float(r.lon) + 180) / 360 * mw
        y = my + (90 - float(r.lat)) / 180 * mh
        if float(r.dP) > 0.01:
            col, rr = CORAL, 5.3
        elif float(r.dP) < -0.01:
            col, rr = BLUE, 5.3
        elif bool(r.fdr_sig):
            col, rr = "#5C91C9", 3.2
        else:
            col, rr = "#7D8895", 2.2
        p.append(circle(x, y, rr, col, "white", 0.8))
    p += [
        badge(mx+30, my+mh-75, "increase-side material", CORAL, "white", 260),
        badge(mx+325, my+mh-75, "depletion-side material", BLUE, "white", 270),
        badge(mx+640, my+mh-75, "detectable direction", "#5C91C9", "white", 260),
    ]

    # Global null gauge + funnel
    p += [panel_label(1135, 185, "B", GREEN), text(1174, 193, "Null first, then regional triage", 29, NAVY, True)]
    gx, gy = 1140, 260
    p += [rect(gx, gy, 455, 310, "#F7FBF6", "#B8D4B4", 2, 24)]
    p += [text(gx+228, gy+62, "GLOBAL MEAN", 23, GREEN, True, anchor="middle"), text(gx+228, gy+124, f"{COUNTS['mean_dp']:+.5f}", 56, NAVY, True, anchor="middle"), text(gx+228, gy+178, "geographic null p = 1.00", 26, GREY, True, anchor="middle")]
    p += [line(gx+75, gy+230, gx+380, gy+230, MID_GREY, 5), circle(gx+228, gy+230, 14, GREEN, "white", 3), text(gx+228, gy+270, "no diffuse global amplification", 22, BLACK, anchor="middle")]
    fx, fy = 1650, 245
    p += [rect(fx, fy, 430, 330, "#F7FAFE", "#B8C8DD", 2, 24)]
    funnel = [("444", "city markers", NAVY), ("311", "detectable", BLUE), ("28", "A/B units", ORANGE), ("21/22", "metro/GHSL", TEAL), ("10", "300-km groups", GREEN)]
    for i, (num, lab, col) in enumerate(funnel):
        y = fy + 52 + i * 57
        p += [circle(fx+62, y, 24, col, "white", 3), text(fx+62, y+8, num, 19, "white", True, anchor="middle"), text(fx+102, y+8, lab, 23, NAVY, True)]
        if i < len(funnel)-1:
            p.append(line(fx+62, y+29, fx+62, y+45, col, 3, marker="arrowBlue"))

    # Regional groups
    p += [panel_label(1135, 640, "C", ORANGE), text(1174, 648, "Regional groups carry the inference", 29, NAVY, True)]
    rx, ry = 1115, 700
    top = REGIONAL.sort_values("n_point_city_units", ascending=False).head(5)
    maxn = top["n_point_city_units"].max()
    for i, (_, r) in enumerate(top.iterrows()):
        y = ry + i * 82
        col = CORAL if r.dominant_direction == "increase-side" else BLUE
        p += [rect(rx, y, 950, 62, "#FFFFFF", "#D3DBE7", 1.5, 15)]
        p += [text(rx+20, y+39, f"Group {int(r.grace_scale_cluster_300km)}", 22, NAVY, True), rect(rx+165, y+16, 500 * r.n_point_city_units / maxn, 28, col, col, 1, 14), text(rx+690, y+39, f"{int(r.n_point_city_units)} cities | {int(r.n_material_units)} material | {r.population_million_sum:.1f}M people", 21, BLACK)]
        p += [text(rx+840, y+39, r.dominant_direction.replace("-side", ""), 20, col, True)]

    p += [
        rect(70, 1240, 2050, 110, "#EEF5FF", "#B8C8DD", 2, 24),
        text(105, 1298, "TAKEAWAY", 28, NAVY, True),
        multiline(270, 1283, "The global null is the guardrail; the useful product is a regional, direction-resolved review queue nested from cities to 300-km GRACE-scale groups.", 25, BLACK, 130),
    ]
    return export_svg("Fig2_global_null_regional_payload_upgrade", p)


def fig3_regional_cards():
    p = svg_header("Figure 3  |  Evidence-Tier Regional Arena", "Five regional cases, five claim classes, one conservative grammar")
    p += [rect(32, 128, 2164, 1270, "#FFFFFF", BORDER, 2.2, 28, True)]
    colors = [GREEN, "#4A9E78", AMBER, BLUE, TEAL]
    for i, (_, r) in enumerate(CARDS.iterrows()):
        x = 80 + i * 420
        y = 230
        p += [rect(x, y, 380, 820, "#FBFCFE", "#BBC9DA", 2, 24, True), panel_label(x+38, y+45, chr(65+i), colors[i])]
        p += [multiline(x+75, y+38, r["regional_unit"], 24, NAVY, 20, bold=True)]
        p += [rect(x+25, y+122, 330, 125, "#F0F5FA", "#D2DCE8", 1.5, 18), multiline(x+43, y+162, r["article_claim_class"], 23, colors[i], 23, bold=True)]
        evs = [
            ("CSR", r["csr_status"]),
            ("GSFC", r["gsfc_status"]),
            ("GFZ", r["gfz_status"]),
            ("Local", r["local_groundwater_or_insar_evidence"]),
        ]
        for j, (lab, val) in enumerate(evs):
            yy = y + 285 + j * 104
            col = colors[i] if j in {0, 3} else BLUE
            p += [circle(x+48, yy+22, 18, col, "white", 2.5), text(x+48, yy+29, str(j+1), 17, "white", True, anchor="middle"), text(x+82, yy+18, lab, 21, NAVY, True), multiline(x+82, yy+45, val, 15, GREY, 37)]
        p += [rect(x+25, y+705, 330, 78, "#FFF8EA" if i == 2 else "#EDF7EF", "#E0C77C" if i == 2 else "#A9CBA6", 1.5, 16), multiline(x+43, y+735, r["claim_boundary"], 15, BLACK, 34)]
    p += [rect(80, 1130, 2040, 190, "#F5F8FC", "#B8C8DD", 2, 24)]
    p += [text(115, 1184, "CLAIM STRENGTH IS NOT A COLORFUL HOTSPOT MAP", 29, NAVY, True)]
    p += [multiline(115, 1240, "Beijing anchors the mechanism, Tokyo Bay/Yokohama is sign-supported, Mumbai-Bhayandar is a contradiction boundary, Delhi is product-material depletion, and Lahore/Punjab is a depletion-subsidence review case.", 25, BLACK, 134)]
    return export_svg("Fig3_regional_evidence_arena_upgrade", p)


def fig4_product_ladder():
    p = svg_header("Figure 4  |  Product-Consensus Guardrail Ladder", "Satellite products and local records tier claims before policy interpretation")
    p += [rect(32, 128, 2164, 1270, "#FFFFFF", BORDER, 2.2, 28, True)]
    p += [panel_label(85, 185, "A", BLUE), text(124, 193, "Mascon and local-evidence ladder", 29, NAVY, True)]
    x0, y0 = 95, 260
    cols = ["regional_card", "CSR", "GSFC", "GFZ raw", "GFZ leakage", "JPL CRI", "Local", "Claim class"]
    widths = [320, 205, 205, 205, 205, 215, 245, 310]
    header_y = y0
    xx = x0
    for c, w in zip(cols, widths):
        p += [rect(xx, header_y, w, 56, NAVY, NAVY, 1, 12), text(xx+w/2, header_y+36, c, 20, "white", True, anchor="middle")]
        xx += w + 8
    for i, (_, r) in enumerate(LADDER.iterrows()):
        y = y0 + 70 + i * 110
        xx = x0
        for c, w in zip(cols, widths):
            val = str(r[c])
            if c == "regional_card":
                fill, col = "#F0F5FA", NAVY
            elif "Material" in val:
                fill, col = BLUE, "white"
            elif "Auth" in val:
                fill, col = "#FCE4D4", NAVY
            elif "Contradicts" in val:
                fill, col = RED, "white"
            elif "High" in val or "Rising" in val or "Borehole" in val:
                fill, col = GREEN, "white"
            elif "Weak" in val or "Near" in val:
                fill, col = "#D9E8F2", NAVY
            else:
                fill, col = SKY, "white"
            p += [rect(xx, y, w, 74, fill, fill if fill != "#F0F5FA" else "#CFD9E6", 1.2, 14)]
            p += [multiline(xx+w/2, y+32, val, 18 if c != "Claim class" else 17, col, max(12, int(w/12)), bold=(c in ["regional_card", "Claim class"]), anchor="middle")]
            xx += w + 8
    # Interpretation framework
    p += [panel_label(85, 905, "B", GREEN), text(124, 913, "Interpretation framework", 29, NAVY, True)]
    boxes = [
        ("SUPPORTED", "Coordinate utility", "Products and local records define a review class.", GREEN),
        ("BOUNDED", "Magnitude materiality", "S_y and local aquifer evidence condition the threshold.", AMBER),
        ("NOT SUPPORTED", "City-scale prediction", "GRACE products do not become wells or engineering design data.", RED),
    ]
    for i, (tag, head, body, col) in enumerate(boxes):
        x = 110 + i * 690
        p += [rect(x, 965, 625, 170, "#F8FAFD", "#D3DBE8", 2, 24), circle(x+55, 1025, 32, col, "white", 4), text(x+55, 1036, "✓" if i < 2 else "×", 31, "white", True, anchor="middle"), text(x+105, 1010, tag, 20, col, True), text(x+105, 1045, head, 25, NAVY, True), multiline(x+105, 1083, body, 20, BLACK, 42)]
    p += [rect(110, 1215, 1980, 92, "#EEF5FF", "#B8C8DD", 2, 22), multiline(145, 1260, "Raw agreement is not enough. The ladder asks: primary screen, independent product sign, leakage stress, authentication boundary, local hydrogeology, then claim class.", 25, BLACK, 135)]
    return export_svg("Fig4_product_evidence_ladder_upgrade", p)


def fig5_sy_phase():
    p = svg_header("Figure 5  |  Aquifer-Class S_y Phase Arena", "Product sign is tested; materiality is aquifer-conditioned")
    p += [rect(32, 128, 2164, 1270, "#FFFFFF", BORDER, 2.2, 28, True)]
    p += [panel_label(82, 185, "A", CORAL), text(121, 193, "S_y prior bands and materiality gate", 29, NAVY, True)]
    x, y = 120, 285
    p += [rect(x, y, 970, 760, "#FBFCFE", "#C7D3E4", 2, 22)]
    names = PHASE["name"].tolist()
    xmax = 0.25
    def sx(v): return x + 250 + v / xmax * 560
    for tick in [0.05, 0.10, 0.15, 0.20, 0.25]:
        p += [line(sx(tick), y+90, sx(tick), y+650, "#D7DFEA", 1.3, "3 7"), text(sx(tick), y+690, f"{tick:.2f}", 16, GREY, anchor="middle")]
    p += [text(x+520, y+730, "specific yield S_y", 21, NAVY, True, anchor="middle")]
    for i, (_, r) in enumerate(PHASE.iterrows()):
        yy = y + 115 + i * 92
        side_col = CORAL if float(r.dP_sy_low) > 0 else BLUE
        p += [text(x+35, yy+7, f"{r['name']} ({r['country']})", 20, NAVY, True), line(sx(r.sy_low), yy, sx(r.sy_high), yy, MID_GREY, 9), circle(sx(r.sy_mid), yy, 12, side_col, "white", 3)]
        p += [text(x+840, yy+7, "material" if bool(r.material_at_sy_mid) else "sub-material", 18, side_col if bool(r.material_at_sy_mid) else GREY, True)]
    p += [line(sx(0.10), y+80, sx(0.10), y+672, NAVY, 2.5, "8 8"), text(sx(0.10)+10, y+75, "baseline S_y=0.10", 18, NAVY, True)]

    # dP phase panel
    p += [panel_label(1190, 185, "B", BLUE), text(1230, 193, "ΔP_liq range under class priors", 29, NAVY, True)]
    x2, y2 = 1190, 285
    p += [rect(x2, y2, 910, 760, "#FBFCFE", "#C7D3E4", 2, 22)]
    minv, maxv = -0.03, 0.03
    def dx(v): return x2 + 430 + (v - minv) / (maxv-minv) * 470
    p += [line(dx(0), y2+90, dx(0), y2+650, "#111827", 2.0), line(dx(-0.01), y2+90, dx(-0.01), y2+650, BLUE, 1.6, "6 6"), line(dx(0.01), y2+90, dx(0.01), y2+650, CORAL, 1.6, "6 6")]
    for tick in [-0.02, -0.01, 0, 0.01, 0.02]:
        p += [text(dx(tick), y2+690, f"{tick:+.2f}", 16, GREY, anchor="middle")]
    p += [text(x2+630, y2+730, "ΔP_liq under class S_y", 21, NAVY, True, anchor="middle")]
    for i, (_, r) in enumerate(PHASE.iterrows()):
        yy = y2 + 115 + i * 92
        vals = [float(r.dP_sy_low), float(r.dP_sy_mid), float(r.dP_sy_high)]
        col = CORAL if max(vals, key=abs) > 0 else BLUE
        p += [text(x2+35, yy+7, r["name"], 20, NAVY, True), line(dx(min(vals)), yy, dx(max(vals)), yy, col, 7), circle(dx(vals[1]), yy, 12, col, "white", 3)]
    # Local replacement rules
    p += [panel_label(1190, 1115, "C", GREEN), text(1230, 1123, "Replacement rule", 28, NAVY, True)]
    rules = [("Sign", "constrained by CSR/GSFC/GFZ"), ("Materiality", "replace S_y prior with local aquifer data"), ("Use", "review cue, not design threshold")]
    for i, (h, b) in enumerate(rules):
        xx = 1240 + i * 300
        p += [rect(xx, 1168, 270, 115, LIGHT_GREEN, "#B7D1B2", 2, 18), text(xx+24, 1210, h, 22, GREEN, True), multiline(xx+24, 1244, b, 17, BLACK, 21)]
    return export_svg("Fig5_aquifer_sy_phase_arena_upgrade", p)


def fig6_engineering_event():
    p = svg_header("Figure 6  |  Engineering Context And Event Boundary", "A stress test for interpretation, not a proof of event-scale prediction")
    p += [rect(32, 128, 2164, 1270, "#FFFFFF", BORDER, 2.2, 28, True)]
    p += [panel_label(85, 185, "A", GREEN), text(124, 193, "Engineering-context enrichment", 29, NAVY, True)]
    x, y = 110, 260
    p += [rect(x, y, 980, 650, "#FBFCFE", "#C7D3E4", 2, 24)]
    top = ENG.head(5)
    for i, (_, r) in enumerate(top.iterrows()):
        yy = y + 90 + i * 105
        frac = min(float(r.followup_fraction_with_proxy), 1)
        cohort = min(float(r.cohort_fraction_with_proxy), 1)
        sig = float(r.fisher_greater_p) < 0.05
        col = GREEN if sig else MID_GREY
        p += [multiline(x+35, yy-12, r.proxy_label, 18, NAVY, 34, bold=True)]
        p += [rect(x+430, yy-23, 400, 24, "#E5EAF1", "#E5EAF1", 1, 12), rect(x+430, yy-23, 400*cohort, 24, "#C7D3E4", "#C7D3E4", 1, 12), rect(x+430, yy+12, 400, 24, "#E5EAF1", "#E5EAF1", 1, 12), rect(x+430, yy+12, 400*frac, 24, col, col, 1, 12)]
        p += [text(x+850, yy-2, f"cohort {cohort:.2f}", 17, GREY), text(x+850, yy+35, f"A/B {frac:.2f}", 17, col, True), text(x+940, yy+17, "✓" if sig else "–", 32, col, True, anchor="middle")]
    p += [badge(x+35, y+590, "enriched proxies support context, not local design", GREEN, "white", 550)]

    p += [panel_label(1180, 185, "B", BLUE), text(1219, 193, "Historical-event boundary benchmark", 29, NAVY, True)]
    x2, y2 = 1180, 260
    p += [rect(x2, y2, 980, 650, "#FBFCFE", "#C7D3E4", 2, 24)]
    events = EVENTS.copy()
    max_abs = max(0.003, events["delta_auc_dynamic_minus_static"].abs().max())
    cx = x2 + 520
    p += [line(cx, y2+90, cx, y2+455, BLACK, 2), text(cx, y2+60, "ΔAUC dynamic - static", 20, NAVY, True, anchor="middle")]
    for i, (_, r) in enumerate(events.iterrows()):
        yy = y2 + 120 + i * 90
        d = float(r.delta_auc_dynamic_minus_static)
        wbar = abs(d)/max_abs*360
        col = GREEN if d > 0 else RED if d < -1e-6 else MID_GREY
        if d >= 0:
            p += [rect(cx, yy-18, wbar, 36, col, col, 1, 8)]
        else:
            p += [rect(cx-wbar, yy-18, wbar, 36, col, col, 1, 8)]
        p += [text(x2+45, yy+8, r.event_title.replace(",", ""), 20, NAVY, True), text(x2+765, yy+8, f"{d:+.4f}", 22, col, True)]
    p += [rect(x2+55, y2+505, 870, 90, "#FFF3F1", "#E8B6AA", 2, 18), multiline(x2+85, y2+542, "Neutral/negative event benchmark delimits the claim: DGLS flags stale water-table assumptions before local modelling; it is not an event-prediction engine.", 22, BLACK, 72)]
    # Protocol strip
    p += [panel_label(85, 990, "C", NAVY), text(124, 998, "Seven-step non-regulatory protocol", 29, NAVY, True)]
    start_x, yy = 120, 1060
    for i, (_, r) in enumerate(PROTOCOL.iterrows()):
        xx = start_x + i * 292
        col = [NAVY, BLUE, CORAL, AMBER, GREEN, TEAL, GREY][i]
        p += [circle(xx+28, yy+35, 30, col, "white", 4), text(xx+28, yy+45, str(int(r.step)), 26, "white", True, anchor="middle"), multiline(xx+70, yy+24, r.review_step, 18, NAVY, 19, bold=True)]
        if i < 6:
            p += [line(xx+205, yy+35, xx+255, yy+35, col, 3, marker="arrowBlue")]
    p += [rect(110, 1245, 1980, 100, "#EEF5FF", "#B8C8DD", 2, 22), multiline(145, 1284, "Use both diagnostics together: engineering context says where local review is plausible; event benchmark prevents overclaiming dynamic event prediction.", 25, BLACK, 132)]
    return export_svg("Fig6_engineering_event_protocol_arena_upgrade", p)


def add_textbox(slide, x, y, w, h, txt, size=18, color=NAVY, bold=False, fill=None, line=None, radius=True):
    if fill:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor.from_string(fill.replace("#", ""))
        shape.line.color.rgb = RGBColor.from_string((line or fill).replace("#", ""))
        tf = shape.text_frame
    else:
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor.from_string(color.replace("#", ""))
    return shape


def add_title(slide, title, subtitle):
    add_textbox(slide, 0.45, 0.18, 12.4, 0.45, title, 24, NAVY, True)
    add_textbox(slide, 0.48, 0.62, 12.0, 0.28, subtitle, 13, NAVY, False)


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slide_defs = [
        ("FIGURE 1 | DYNAMIC GROUNDWATER REVIEW SCREEN", "Editable 3D mechanism storyboard", ["GRACE regional storage", "3D aquifer cutaway", "Zhu water-table term", "A/B/C/D review action"]),
        ("FIGURE 2 | GLOBAL NULL, REGIONAL PAYLOAD", "Counts and regional hierarchy", ["444 cities", "311 detectable", "28 A/B units", "10 regional groups"]),
        ("FIGURE 3 | EVIDENCE-TIER REGIONAL ARENA", "Five evidence classes", list(CARDS["regional_unit"].head(5))),
        ("FIGURE 4 | PRODUCT-CONSENSUS GUARDRAIL LADDER", "CSR / GSFC / GFZ / JPL / Local", list(LADDER["regional_card"])),
        ("FIGURE 5 | AQUIFER-CLASS S_y PHASE ARENA", "S_y controls materiality", list(PHASE["name"])),
        ("FIGURE 6 | ENGINEERING CONTEXT AND EVENT BOUNDARY", "Context support plus event-claim boundary", ["engineering proxies", "event ΔAUC", "protocol steps", "non-regulatory use"]),
        ("TABLE/BOX 1 | REVIEW PROTOCOL", "Editable protocol schematic", list(PROTOCOL["review_step"])),
    ]
    for idx, (title, subtitle, items) in enumerate(slide_defs):
        slide = prs.slides.add_slide(blank)
        add_title(slide, title, subtitle)
        add_textbox(slide, 0.35, 1.05, 12.6, 5.95, "", 12, NAVY, False, "#F8FAFE", "#B8C8DD")
        if idx == 0:
            # isometric aquifer and arrows
            add_textbox(slide, 0.75, 1.35, 3.2, 1.25, "3D aquifer\nwater-table state", 20, "FFFFFF", True, BLUE)
            add_textbox(slide, 4.45, 1.35, 2.4, 1.0, "GRACE\nregional storage", 18, "FFFFFF", True, TEAL)
            add_textbox(slide, 7.35, 1.35, 2.3, 1.0, "Zhu model\nWTD term", 18, "FFFFFF", True, CORAL)
            add_textbox(slide, 10.1, 1.35, 2.1, 1.0, "local review\naction", 18, "FFFFFF", True, GREEN)
            add_textbox(slide, 0.9, 4.1, 11.4, 1.2, "recharge -> shallower WTD -> P_liq up     |     depletion -> deeper WTD -> P_liq down but subsidence/water-security audit", 18, NAVY, True, "#EEF5FF", "#B8C8DD")
        elif idx == 1:
            for i, (num, lab, col) in enumerate([("444", "city markers", NAVY), ("311", "detectable", BLUE), ("28", "A/B units", ORANGE), ("21/22", "metro/GHSL", TEAL), ("10", "regional groups", GREEN)]):
                add_textbox(slide, 0.9+i*2.45, 1.65, 1.4, 0.7, num, 24, "FFFFFF", True, col)
                add_textbox(slide, 0.72+i*2.45, 2.45, 1.8, 0.45, lab, 13, NAVY, True)
            add_textbox(slide, 1.0, 4.0, 11.4, 1.0, "Global mean near zero; geographic null p = 1.00. The useful result is regional and bidirectional.", 19, NAVY, True, "#EAF5E8", "#B8D4B4")
        elif idx in {2, 3}:
            for i, item in enumerate(items[:5]):
                add_textbox(slide, 0.72+i*2.48, 1.45, 2.15, 3.75, str(item), 15, NAVY, True, "#FFFFFF", "#B8C8DD")
                add_textbox(slide, 0.88+i*2.48, 4.9, 1.82, 0.5, ["anchor", "sign", "candidate", "material", "review"][i] if idx == 2 else ["CSR", "GSFC", "GFZ", "JPL", "Local"][i], 13, "FFFFFF", True, [GREEN, TEAL, AMBER, BLUE, CORAL][i])
        elif idx == 4:
            for i, item in enumerate(items):
                add_textbox(slide, 0.7, 1.25+i*0.75, 2.0, 0.38, str(item), 12, NAVY, True)
                add_textbox(slide, 2.8, 1.28+i*0.75, 6.0, 0.22, "", 8, NAVY, False, "#D7DEE8", "#D7DEE8", radius=False)
                add_textbox(slide, 5.25, 1.22+i*0.75, 0.35, 0.35, "", 8, "FFFFFF", False, CORAL if i < 3 else BLUE)
            add_textbox(slide, 9.2, 2.0, 2.7, 2.5, "Replacement rule\nSign: products\nMateriality: local S_y\nUse: review cue", 16, NAVY, True, "#EEF5FF", "#B8C8DD")
        elif idx == 5:
            for i, item in enumerate(items):
                add_textbox(slide, 0.9+i*3.0, 1.55, 2.4, 1.1, item, 17, NAVY, True, "#FFFFFF", "#B8C8DD")
            add_textbox(slide, 1.0, 4.15, 11.2, 1.0, "Event benchmark is neutral/negative for broad dynamic superiority: use DGLS for pre-implementation review, not event prediction.", 18, RED, True, "#FFF3F1", "#E8B6AA")
        else:
            for i, item in enumerate(items[:7]):
                add_textbox(slide, 0.65+i*1.8, 1.45, 1.45, 0.7, str(i+1), 20, "FFFFFF", True, [NAVY, BLUE, CORAL, AMBER, GREEN, TEAL, GREY][i])
                add_textbox(slide, 0.45+i*1.8, 2.25, 1.75, 1.1, item, 12, NAVY, True, "#FFFFFF", "#B8C8DD")
        add_textbox(slide, 0.65, 6.62, 12.0, 0.36, "Editable storyboard: shapes and text can be modified in PowerPoint; final publication figures are vector SVG/PDF.", 10, GREY)
    prs.save(PPT_DIR / "NatureWater_Editable_Figure_Upgrade_Pack.pptx")


def build_contact_sheet(figs: list[Path]) -> Path:
    from PIL import Image, ImageDraw

    pngs = [PNG_DIR / f"{p.stem}.png" for p in figs if (PNG_DIR / f"{p.stem}.png").exists()]
    thumb_w, thumb_h = 720, 491
    pad = 32
    label_h = 36
    rows = math.ceil(len(pngs) / 2)
    sheet = Image.new("RGB", (2 * thumb_w + 3 * pad, rows * (thumb_h + label_h) + (rows + 1) * pad), "white")
    draw = ImageDraw.Draw(sheet)
    for i, png in enumerate(pngs):
        img = Image.open(png).convert("RGB")
        img.thumbnail((thumb_w, thumb_h), Image.LANCZOS)
        col, row = i % 2, i // 2
        x = pad + col * (thumb_w + pad)
        y = pad + row * (thumb_h + label_h + pad)
        draw.text((x, y), png.stem, fill=(8, 27, 74))
        sheet.paste(img, (x, y + label_h))
    out = PNG_DIR / "contact_sheet.png"
    sheet.save(out, quality=95)
    return out


def box1_protocol():
    p = svg_header("Box 1  |  Dynamic Groundwater-Liquefaction Review Protocol", "A non-regulatory local-data replacement workflow")
    p += [rect(32, 128, 2164, 1270, "#FFFFFF", BORDER, 2.2, 28, True)]
    colors = [NAVY, BLUE, CORAL, AMBER, GREEN, TEAL, GREY]
    start_x, start_y = 105, 270
    for i, r in enumerate(PROTOCOL.itertuples()):
        x = start_x + (i % 4) * 510
        y = start_y + (i // 4) * 380
        col = colors[i % len(colors)]
        step_text = getattr(r, "review_step", f"Step {i+1}")
        required = getattr(r, "required_evidence", "required local evidence")
        output = getattr(r, "decision_output", "review output")
        p += [
            rect(x, y, 455, 270, "#FBFCFE", "#B8C8DD", 2, 22, True),
            circle(x + 54, y + 58, 34, col, "white", 4),
            text(x + 54, y + 70, str(int(getattr(r, "step", i + 1))), 28, "white", True, anchor="middle"),
            multiline(x + 105, y + 52, step_text, 25, NAVY, 24, bold=True),
            multiline(x + 38, y + 128, required, 17, GREY, 42),
            rect(x + 32, y + 218, 390, 38, "#EEF5FF", "#D3DBE7", 1.2, 13),
            multiline(x + 46, y + 242, output, 15, BLACK, 48),
        ]
        if i < len(PROTOCOL) - 1 and i % 4 != 3:
            p += [line(x + 455, y + 134, x + 500, y + 134, col, 4, marker="arrowBlue")]
    p += [
        rect(110, 1110, 1980, 170, "#EEF5FF", "#B8C8DD", 2, 24),
        text(145, 1170, "Outcome classes", 30, NAVY, True),
        badge(430, 1138, "routine monitoring", GREY, "white", 310),
        badge(775, 1138, "targeted data collection", ORANGE, "white", 380),
        badge(1195, 1138, "local geotechnical review", BLUE, "white", 390),
        badge(1625, 1138, "multi-hazard audit", GREEN, "white", 310),
        multiline(145, 1232, "Protocol output is a review action. It does not create a regulatory threshold, hazard map, event-prediction model, or factor of safety.", 24, BLACK, 126),
    ]
    return export_svg("Box1_review_protocol", p)


FIGURE_BUILDERS = {
    "Fig1": fig1_mechanism,
    "Fig2": fig2_global_payload,
    "Fig3": fig3_regional_cards,
    "Fig4": fig4_product_ladder,
    "Fig5": fig5_sy_phase,
    "Fig6": fig6_engineering_event,
    "Box1": box1_protocol,
}


def main(argv: list[str] | None = None):
    parser = argparse.ArgumentParser()
    parser.add_argument("--fig", choices=sorted(FIGURE_BUILDERS), help="Build one figure only.")
    parser.add_argument("--no-pptx", action="store_true", help="Skip editable PPTX deck.")
    args = parser.parse_args(argv)
    figs = [FIGURE_BUILDERS[args.fig]()] if args.fig else [builder() for builder in FIGURE_BUILDERS.values()]
    if not args.no_pptx:
        build_pptx()
    build_contact_sheet(figs)
    review = ROOT / "figure_specs" / "visual_upgrade_pack_review.md"
    review.parent.mkdir(parents=True, exist_ok=True)
    review.write_text(
        "# Nature Water Visual Upgrade Pack Review\n\n"
        "Generated SVG/PDF/PNG candidates for Figures 1-6 and Box 1 plus an editable PPTX deck.\n\n"
        "Scientific boundary: figures explain the water-table state-variable pathway and evidence tiers only; no AI-generated raster is used as evidence.\n\n"
        "Self-review: PASS for editable vector provenance; final manual typography micro-adjustment remains optional.\n\n"
        "Outputs:\n" + "\n".join(f"- `{p.name}`" for p in figs) + "\n",
        encoding="utf-8",
    )


if __name__ == "__main__":
    main()
